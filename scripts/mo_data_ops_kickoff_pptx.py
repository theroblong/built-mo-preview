"""
Generate Mo Data Operations Kickoff presentation (.pptx)
Output: mockups/mo_data_ops_kickoff.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1B, 0x35, 0x66)
NAVY_DEEP   = RGBColor(0x0C, 0x1A, 0x38)
BLUE        = RGBColor(0x40, 0x78, 0xD0)
BLUE_LIGHT  = RGBColor(0xA0, 0xC0, 0xF0)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF0, 0xF4, 0xFA)
LIGHT_BLUE  = RGBColor(0xD8, 0xE5, 0xFF)
MID_GREY    = RGBColor(0x6B, 0x7F, 0xA0)
DARK_TEXT   = RGBColor(0x11, 0x18, 0x27)
GREEN       = RGBColor(0x0E, 0x5E, 0x38)
GREEN_LIGHT = RGBColor(0xD4, 0xF0, 0xE4)
AMBER       = RGBColor(0x8A, 0x48, 0x00)
AMBER_LIGHT = RGBColor(0xFF, 0xF3, 0xD6)
PURPLE      = RGBColor(0x4B, 0x2D, 0x8A)
PURPLE_LIGHT= RGBColor(0xEC, 0xE8, 0xF8)
RED         = RGBColor(0x8B, 0x1A, 0x1A)
RED_LIGHT   = RGBColor(0xFE, 0xF0, 0xF0)
GREY_MID    = RGBColor(0x4A, 0x56, 0x6A)
GREY_LIGHT  = RGBColor(0xD4, 0xDA, 0xE6)
GREY_PALE   = RGBColor(0xEE, 0xF0, 0xF4)

# ── Slide dimensions (widescreen 16:9) ───────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────
def tb(slide, left, top, width, height, text, size=12, bold=False, color=DARK_TEXT,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def tb_lines(slide, left, top, width, height, lines, default_size=11,
             default_color=DARK_TEXT, default_bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, bold, color, indent = item, default_size, default_bold, default_color, False
        else:
            text  = item[0]
            size  = item[1] if len(item) > 1 else default_size
            bold  = item[2] if len(item) > 2 else default_bold
            color = item[3] if len(item) > 3 else default_color
            indent= item[4] if len(item) > 4 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if indent:
            p.level = 1
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def rect(slide, left, top, width, height, fill=NAVY, line=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.width = Pt(line)
        if line_color:
            shape.line.color.rgb = line_color
    return shape


def footnote(slide):
    tb(slide, 0.5, 7.2, 12.5, 0.25,
       "Mo by Aevah  ·  BUILT Kickoff  ·  August 2026  ·  Confidential",
       size=8, color=MID_GREY, align=PP_ALIGN.RIGHT)


def phase_badge(slide, phase):
    """Top-right phase pill on module slides."""
    if phase == 1:
        bg, text_color, label = BLUE, WHITE, "PHASE 1  ·  Now – Dec 2026"
    else:
        bg, text_color, label = GREY_MID, WHITE, "PHASE 2  ·  Planned 2027"
    rect(slide, 9.8, 0.22, 3.3, 0.32, fill=bg)
    tb(slide, 9.85, 0.26, 3.2, 0.24, label,
       size=8.5, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def raci_strip(slide, aevah_text, built_text):
    """Two-column RACI strip at the bottom of each module slide."""
    y = 6.12
    h = 0.88
    left_w = 5.0
    right_x = 5.85
    right_w = 7.1

    # Aevah side
    rect(slide, 0.45, y, left_w, 0.26, fill=NAVY)
    tb(slide, 0.55, y + 0.04, left_w - 0.2, 0.2,
       "AEVAH DELIVERS", size=7.5, bold=True, color=BLUE_LIGHT)
    rect(slide, 0.45, y + 0.26, left_w, h - 0.26,
         fill=RGBColor(0xE8, 0xED, 0xF8), line=0.5, line_color=RGBColor(0xB0, 0xC0, 0xE0))
    tb(slide, 0.55, y + 0.32, left_w - 0.2, h - 0.32,
       aevah_text, size=8.5, color=DARK_TEXT)

    # BUILT side
    rect(slide, right_x, y, right_w, 0.26, fill=GREEN)
    tb(slide, right_x + 0.1, y + 0.04, right_w - 0.2, 0.2,
       "BUILT PROVIDES", size=7.5, bold=True, color=WHITE)
    rect(slide, right_x, y + 0.26, right_w, h - 0.26,
         fill=GREEN_LIGHT, line=0.5, line_color=RGBColor(0x80, 0xC8, 0xA0))
    tb(slide, right_x + 0.1, y + 0.32, right_w - 0.2, h - 0.32,
       built_text, size=8.5, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.333, 7.5, fill=NAVY_DEEP)
rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)

tb(slide, 0.5, 0.8, 10, 0.4,
   "MO BY AEVAH  ·  BUILT KICKOFF  ·  AUGUST 2026",
   size=9, color=BLUE_LIGHT, bold=True)

tb(slide, 0.5, 1.35, 10, 1.8,
   "How Mo Works\nWith Your Data",
   size=42, bold=True, color=WHITE)

tb(slide, 0.5, 3.35, 8.5, 1.2,
   "A walkthrough of the process behind each Mo intelligence module — "
   "what data it needs, where that data comes from, what operations are performed, "
   "and what appears on screen.",
   size=14, color=BLUE_LIGHT)

phase1_modules = ["Cannibalization Risk", "Price Elasticity"]
phase2_modules = ["Promotional Response", "Demand Forecast", "Launch Monitoring"]

# Phase 1 — header band above module boxes
p1_y = 4.3
rect(slide, 0.5, p1_y, 12.4, 0.34, fill=RGBColor(0x1A, 0x34, 0x62))
tb(slide, 0.65, p1_y + 0.07, 11.5, 0.22,
   "PHASE 1  ·  NOW THROUGH DEC 2026", size=11, bold=True, color=BLUE_LIGHT)
for i, mod in enumerate(phase1_modules):
    x = 0.5 + i * 6.3
    rect(slide, x, p1_y + 0.34, 6.1, 0.75, fill=RGBColor(0x20, 0x40, 0x70))
    tb(slide, x + 0.15, p1_y + 0.57, 5.8, 0.38,
       mod, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Phase 2 — header band above module boxes
p2_y = p1_y + 0.34 + 0.75 + 0.2
rect(slide, 0.5, p2_y, 12.4, 0.34, fill=RGBColor(0x28, 0x32, 0x44))
tb(slide, 0.65, p2_y + 0.07, 11.5, 0.22,
   "PHASE 2  ·  PLANNED 2027", size=11, bold=True, color=GREY_LIGHT)
for i, mod in enumerate(phase2_modules):
    x = 0.5 + i * 4.2
    rect(slide, x, p2_y + 0.34, 4.0, 0.65, fill=RGBColor(0x2E, 0x38, 0x4A))
    tb(slide, x + 0.1, p2_y + 0.52, 3.8, 0.34,
       mod, size=12, bold=True, color=GREY_LIGHT, align=PP_ALIGN.CENTER)

tb(slide, 9.0, 7.1, 4, 0.3,
   "aevah.ai  ·  Confidential",
   size=9, color=MID_GREY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PIPELINE OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.333, 7.5, fill=OFF_WHITE)
rect(slide, 0, 0, 13.333, 1.1, fill=NAVY)
rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)

tb(slide, 0.5, 0.22, 10, 0.6,
   "The End-to-End Process",
   size=24, bold=True, color=WHITE)

tb(slide, 0.5, 1.2, 12.5, 0.4,
   "Five stages from your SPINS export to decisions on screen — Aevah automates stages 2–4 so no manual steps are required once data is deposited.",
   size=11.5, color=MID_GREY)

steps = [
    ("Stage 1", "SPINS Export", "BUILT Data Team",
     "214-column weekly POS extract deposited to shared storage"),
    ("Stage 2", "Data Ingest", "Aevah — automated",
     "File-watch trigger detects new export; quality gates applied; incremental load only"),
    ("Stage 3", "Feature\nEngineering", "Aevah — automated",
     "Aevah compresses raw rows into ML-ready feature tables per module — no manual steps"),
    ("Stage 4", "Model Training\n& Scoring", "Aevah — automated",
     "Models retrain on drift/schedule; scores saved to Aevah; smoke-test gates at each stage"),
    ("Stage 5", "Mo UI", "BUILT Team",
     "Two data layers: ML projections (forecast) + live context data (actuals at render time)"),
]

step_w = 2.2
box_h  = 2.85
for i, (num, name, who, desc) in enumerate(steps):
    x = 0.4 + i * (step_w + 0.2)
    y = 1.88

    rect(slide, x, y, step_w, box_h, fill=NAVY)
    tb(slide, x + 0.12, y + 0.1, step_w - 0.24, 0.28,
       num.upper(), size=8, bold=True, color=BLUE_LIGHT)
    tb(slide, x + 0.12, y + 0.38, step_w - 0.24, 0.55,
       name, size=13, bold=True, color=WHITE)
    tb(slide, x + 0.12, y + 0.92, step_w - 0.24, 0.3,
       who, size=8.5, color=BLUE_LIGHT, bold=True)
    tb(slide, x + 0.12, y + 1.25, step_w - 0.24, 1.4,
       desc, size=9, color=RGBColor(0xB0, 0xC8, 0xF0))

    if i < len(steps) - 1:
        ax = x + step_w + 0.02
        tb(slide, ax, y + 1.2, 0.2, 0.35, "→", size=18, bold=True,
           color=BLUE, align=PP_ALIGN.CENTER)

rect(slide, 0.4, 5.1, 12.5, 1.15, fill=LIGHT_BLUE)
tb(slide, 0.6, 5.17, 12, 0.28, "KEY PRINCIPLE", size=8, bold=True, color=BLUE)
tb(slide, 0.6, 5.45, 12, 0.7,
   "Every Mo screen draws from two distinct sources: ML scoring tables (what the model forecasts, written in Stage 4) "
   "and live Aevah queries (actual component data queried at render time). "
   "These have separate lineage and must never be blended.",
   size=10.5, color=DARK_TEXT)

footnote(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ROLES & RESPONSIBILITIES (RACI)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.333, 7.5, fill=OFF_WHITE)
rect(slide, 0, 0, 13.333, 1.1, fill=NAVY)
rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)

tb(slide, 0.5, 0.18, 10, 0.6,
   "Roles & Responsibilities — Phase 1",
   size=24, bold=True, color=WHITE)
tb(slide, 0.5, 0.72, 12.5, 0.32,
   "What Aevah commits to deliver, what BUILT commits to provide, and the shared milestones where both teams align.",
   size=11, color=BLUE_LIGHT)

# ── Left column: Aevah Delivers ──
left_x, left_w = 0.42, 6.0
rect(slide, left_x, 1.22, left_w, 0.32, fill=NAVY)
tb(slide, left_x + 0.12, 1.26, left_w - 0.24, 0.24,
   "AEVAH DELIVERS", size=9, bold=True, color=BLUE_LIGHT)

aevah_items = [
    "Platform setup, hosting, and data ingestion pipeline",
    "Data quality gates at every stage — no bad data reaches the UI",
    "Feature engineering: raw SPINS rows compressed into ML-ready tables per module",
    "Model training and scoring: cannibalization risk + price elasticity (Phase 1)",
    "Mo UI deployment: all screens, filters, and visualizations",
    "Model monitoring: drift detection and automatic retraining on schedule",
    "Customer-facing overview and reference documentation, updated each session",
]
item_y = 1.54
for item in aevah_items:
    rect(slide, left_x + 0.12, item_y + 0.12, 0.08, 0.08, fill=BLUE)
    tb(slide, left_x + 0.3, item_y + 0.04, left_w - 0.42, 0.32,
       item, size=10, color=DARK_TEXT)
    item_y += 0.38

# ── Right column: BUILT Provides ──
right_x, right_w = 6.85, 6.1
rect(slide, right_x, 1.22, right_w, 0.32, fill=GREEN)
tb(slide, right_x + 0.12, 1.26, right_w - 0.24, 0.24,
   "BUILT PROVIDES", size=9, bold=True, color=WHITE)

built_items = [
    "Weekly SPINS export (214 columns) deposited to shared storage on agreed schedule",
    "ERP / data warehouse access — for Phase 2 data sources and Connor's actuals",
    "Stakeholder availability: kickoff, data QA review, demo sessions, and sign-off",
    "Projection validation: Connor + Brian confirm whether findings match their market knowledge",
    "UI feedback: what information helps the job get done — what to add, change, or remove",
    "Phase 1 milestone sign-off before proceeding to Phase 2 scope",
]
item_y = 1.54
for item in built_items:
    rect(slide, right_x + 0.12, item_y + 0.12, 0.08, 0.08, fill=GREEN)
    tb(slide, right_x + 0.3, item_y + 0.04, right_w - 0.42, 0.32,
       item, size=10, color=DARK_TEXT)
    item_y += 0.38

# ── RACI Matrix ──
matrix_y = 4.38
col_widths = [4.2, 1.8, 1.5, 1.6, 1.8, 2.0]  # Activity, Aevah, Brian, Connor, Data Team, Milestone
col_x = [0.42]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)

headers = ["Activity", "Aevah", "Brian", "Connor", "BUILT Data", "Milestone"]
rect(slide, 0.42, matrix_y, 12.53, 0.3, fill=NAVY)
for i, (h, x) in enumerate(zip(headers, col_x)):
    align = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
    pad = 0.1 if i == 0 else 0.0
    tb(slide, x + pad, matrix_y + 0.05, col_widths[i] - 0.1, 0.22,
       h, size=8, bold=True, color=WHITE, align=align)
matrix_y += 0.3

raci_rows = [
    ("SPINS export & weekly deposit to shared storage",
     "I", "A", "—", "R", "Week 1–2"),
    ("Data ingest, quality gates, feature engineering",
     "R", "I", "—", "—", "Week 3–4"),
    ("Model training, scoring & Mo UI deployment",
     "R", "I", "—", "—", "Week 4–6"),
    ("Projection review & business validation",
     "C", "A", "R", "—", "Week 6–8"),
    ("UI feedback sessions",
     "C", "A", "R", "—", "Week 7–10"),
    ("Phase 1 milestone sign-off",
     "A", "A", "C", "—", "Week 10–12"),
]

for row_i, (activity, *cells) in enumerate(raci_rows):
    row_bg = RGBColor(0xF8, 0xFA, 0xFD) if row_i % 2 == 0 else WHITE
    row_h = 0.32
    rect(slide, 0.42, matrix_y, 12.53, row_h, fill=row_bg,
         line=0.5, line_color=RGBColor(0xCC, 0xD8, 0xEE))
    tb(slide, col_x[0] + 0.1, matrix_y + 0.06, col_widths[0] - 0.2, row_h - 0.08,
       activity, size=8.5, color=DARK_TEXT)
    for i, (cell, cx) in enumerate(zip(cells, col_x[1:])):
        cell_color = BLUE if cell == "R" else (GREEN if cell == "A" else DARK_TEXT)
        tb(slide, cx, matrix_y + 0.06, col_widths[i + 1], row_h - 0.08,
           cell, size=9, bold=(cell in ("R", "A")), color=cell_color,
           align=PP_ALIGN.CENTER)
    matrix_y += row_h

# Legend
tb(slide, 0.42, matrix_y + 0.08, 8, 0.24,
   "R = Responsible (does the work)   A = Accountable (owns the outcome)   C = Consulted   — = Not in scope for this activity",
   size=7.5, color=MID_GREY, italic=True)

footnote(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE TEMPLATE FUNCTION — simulation modules
# ═══════════════════════════════════════════════════════════════════════════════
def add_sim_slide(num_label, name, answer, inputs, operations, ml_outputs, live_outputs,
                  phase=1, aevah_raci="", built_raci=""):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.333, 7.5, fill=OFF_WHITE)
    rect(slide, 0, 0, 13.333, 1.4, fill=NAVY)
    rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)

    # Phase 2 muted overlay banner
    if phase == 2:
        rect(slide, 0, 0, 13.333, 1.4, fill=GREY_MID)

    phase_str = "PHASE 1 — NOW THROUGH DEC 2026" if phase == 1 else "PHASE 2 — PLANNED 2027"
    phase_col = BLUE if phase == 1 else GREY_LIGHT
    tb(slide, 0.5, 0.1, 3.5, 0.28, num_label, size=9, bold=True, color=BLUE_LIGHT)
    tb(slide, 4.0, 0.1, 9.0, 0.28, phase_str, size=9, bold=True, color=phase_col,
       align=PP_ALIGN.RIGHT)
    tb(slide, 0.5, 0.38, 7, 0.6, name, size=22, bold=True, color=WHITE)
    tb(slide, 7.8, 0.38, 5.3, 0.8, f'"{answer}"',
       size=11, color=RGBColor(0xC0, 0xD5, 0xFF), italic=True)

    left_x = 0.45
    left_w = 5.0
    y = 1.55

    rect(slide, left_x, y, left_w, 0.28, fill=NAVY)
    tb(slide, left_x + 0.1, y + 0.04, left_w - 0.2, 0.22,
       "DATA REQUIRED", size=8, bold=True, color=BLUE_LIGHT)
    y += 0.28

    rect(slide, left_x, y, left_w, 0.27, fill=LIGHT_BLUE)
    tb(slide, left_x + 0.1, y + 0.04, 2.4, 0.2, "Field", size=8, bold=True, color=DARK_TEXT)
    tb(slide, left_x + 2.6, y + 0.04, 1.0, 0.2, "Source", size=8, bold=True, color=DARK_TEXT)
    tb(slide, left_x + 3.7, y + 0.04, 1.2, 0.2, "Purpose", size=8, bold=True, color=DARK_TEXT)
    y += 0.27

    row_h = 0.31
    for i, (field, source, purpose) in enumerate(inputs):
        bg = RGBColor(0xF8, 0xFA, 0xFD) if i % 2 == 0 else WHITE
        rect(slide, left_x, y, left_w, row_h, fill=bg, line=0.5,
             line_color=RGBColor(0xCC, 0xD8, 0xEE))
        tb(slide, left_x + 0.08, y + 0.06, 2.4, row_h - 0.06,
           field, size=8.5, bold=True, color=DARK_TEXT)
        tb(slide, left_x + 2.58, y + 0.06, 1.0, row_h - 0.06,
           source, size=8, color=BLUE, bold=True)
        tb(slide, left_x + 3.68, y + 0.06, 1.22, row_h - 0.06,
           purpose, size=7, color=MID_GREY)
        y += row_h

    y += 0.12
    rect(slide, left_x, y, left_w, 0.28, fill=NAVY)
    tb(slide, left_x + 0.1, y + 0.04, left_w - 0.2, 0.22,
       "OPERATIONS PERFORMED", size=8, bold=True, color=BLUE_LIGHT)
    y += 0.28

    for op in operations:
        rect(slide, left_x + 0.1, y + 0.1, 0.08, 0.08, fill=BLUE)
        tb(slide, left_x + 0.28, y + 0.02, left_w - 0.38, 0.3,
           op, size=8.5, color=DARK_TEXT)
        y += 0.3

    right_x = 5.85
    right_w  = 7.1
    y_r = 1.55

    rect(slide, right_x, y_r, right_w, 0.3, fill=GREEN)
    tb(slide, right_x + 0.12, y_r + 0.05, right_w, 0.22,
       "FORECAST & PROJECTION — from trained model (Stage 4)", size=8, bold=True, color=WHITE)
    y_r += 0.3

    rect(slide, right_x, y_r, right_w, len(ml_outputs) * 0.35 + 0.1,
         fill=GREEN_LIGHT, line=0.5, line_color=RGBColor(0x68, 0xC9, 0xA0))
    for ml_out in ml_outputs:
        tb(slide, right_x + 0.3, y_r + 0.07, right_w - 0.4, 0.28,
           f"›  {ml_out}", size=9.5, color=GREEN)
        y_r += 0.35
    y_r += 0.1

    y_r += 0.12
    rect(slide, right_x, y_r, right_w, 0.3, fill=AMBER)
    tb(slide, right_x + 0.12, y_r + 0.05, right_w, 0.22,
       "LIVE CONTEXT DATA — queried at display time (Stage 5)", size=8, bold=True, color=WHITE)
    y_r += 0.3

    rect(slide, right_x, y_r, right_w, len(live_outputs) * 0.35 + 0.1,
         fill=AMBER_LIGHT, line=0.5, line_color=RGBColor(0xE8, 0xB8, 0x70))
    for lv_out in live_outputs:
        tb(slide, right_x + 0.3, y_r + 0.07, right_w - 0.4, 0.28,
           f"›  {lv_out}", size=9.5, color=AMBER)
        y_r += 0.35

    # Per-module RACI strip
    if aevah_raci or built_raci:
        raci_strip(slide, aevah_raci, built_raci)

    footnote(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES 4–8 — Simulation modules
# ═══════════════════════════════════════════════════════════════════════════════

add_sim_slide(
    num_label="MODULE 01",
    name="Cannibalization Risk",
    answer="Is this new SKU pulling demand from an existing one, or adding net-new buyers?",
    inputs=[
        ("Weekly units sold", "SPINS", "Before/after demand"),
        ("Base units (non-promo)", "SPINS", "Isolate organic demand"),
        ("Distribution — stores selling", "SPINS", "Demand vs. distrib. loss"),
        ("Average retail price", "SPINS", "Rule out price effect"),
        ("First week selling (launch date)", "SPINS", "Launch anchor date"),
        ("Pack count, flavor, brand line", "SPINS", "Find substitute SKUs"),
        ("Weeks on promotion", "SPINS", "Flag promo periods"),
    ],
    operations=[
        "Build before/after windows (4-week, 13-week, 26-week) around each product launch",
        "Pair each focal SKU with candidate donors — same flavor, same pack family, competitor",
        "Measure whether focal demand rose while donor demand fell in the same stores and window",
        "Flag windows contaminated by promotions, distribution ramps, or supply events",
        "Calculate incremental share: net-new demand vs. demand transferred from a donor SKU",
        "Assign label (Cannibalizing / Watch / Incremental / Neutral) from observable evidence",
    ],
    ml_outputs=[
        "Risk score + confidence for each focal/donor pair",
        "Status: Cannibalizing · Watch · Incremental · Neutral",
        "Ranked donors with evidence summary (focal lift, donor decline, distribution delta)",
        "Incremental share — net-new demand vs. transferred demand  [FML-10]",
        "Cannibalization rate actuals + 13-week forward projection  [FML-09]",
    ],
    live_outputs=[
        "Actual weekly demand trend for focal and donor SKUs",
        "Distribution (TDP) trend — stores stocking each SKU over time",
        "Geography heatmap — which markets show the strongest demand transfer",
        "Launch ramp chart — weeks 1–16 of distribution build vs. category norms",
    ],
    phase=1,
    aevah_raci="Comparison pool build · Pairing model · Risk scoring · Publish to Mo UI · Monitor model drift",
    built_raci="SPINS export deposit · Review risk labels — do the identified donor pairs make business sense? · Validate threshold cutoffs with Brian",
)

add_sim_slide(
    num_label="MODULE 02",
    name="Price Elasticity",
    answer="If we change price by X%, how much does demand change — and what is the optimal price point?",
    inputs=[
        ("Average retail price", "SPINS", "Price level measure"),
        ("Weekly units sold", "SPINS", "Demand response"),
        ("Base price (non-promo ARP)", "SPINS", "Everyday vs. promo"),
        ("Promo discount depth (%)", "SPINS", "Quantify promo depth"),
        ("Incremental units (promo-driven)", "SPINS", "Isolate promo demand"),
        ("Distribution (TDP)", "SPINS", "Distribution control"),
        ("Competitor price (same flavor)", "SPINS", "Own vs. competitor"),
    ],
    operations=[
        "Normalize price to $/bar so singles and 12-packs are comparable",
        "Build 13-week log-log regression windows per SKU × account",
        "Separate promotional price changes from everyday price changes",
        "Compare BUILT $/bar vs. same-flavor competitor in the same account and week",
        "Compare $/bar across BUILT own pack sizes of the same flavor (pack-ladder check)",
        "Reject windows where price moved >40% in one week — data artifact guardrail",
    ],
    ml_outputs=[
        "Own-price elasticity ε with confidence range (q10 / q50 / q90)",
        "Cross-price elasticity — demand response to competitor price changes",
        "Promo elasticity curve — lift at each discount depth",
        "What-if demand curves — projected units at alternative price points",
        "Pricing event queue — active competitive gaps or pack-ladder compression",
    ],
    live_outputs=[
        "Current $/bar and 52-week price trend by retailer account",
        "Pack price ladder — single / 4-pack / 12-pack side-by-side in the same account",
        "Competitive price gap chart — BUILT vs. Tier 1 competitor in the same flavor",
        "Category benchmark — BUILT $/bar vs. MULO FOOD category average",
    ],
    phase=1,
    aevah_raci="$/bar normalization · Elasticity regression models · Score publishing · Pack-ladder + competitor comparisons · Publish to Mo UI",
    built_raci="SPINS export deposit · Validate elasticity values — do they match BUILT's commercial experience? · Confirm price-point recommendations with Brian",
)

add_sim_slide(
    num_label="MODULE 03",
    name="Promotional Response",
    answer="How much incremental demand does a promotion generate — and what depth maximizes return?",
    inputs=[
        ("Incremental units", "SPINS", "Above-baseline demand"),
        ("Base units (non-promo)", "SPINS", "Lift denominator"),
        ("% units sold on promotion", "SPINS", "Promo dependency"),
        ("Promo discount depth (%)", "SPINS", "Quantify discount"),
        ("Channel (grocery / mass / c-store)", "SPINS", "Channel response"),
        ("Pack count", "SPINS", "Pack size behavior"),
        ("Distribution on promo (TDP, Any Promo)", "SPINS", "Display vs. price"),
    ],
    operations=[
        "Calculate promo lift ratio: incremental units ÷ base units, per SKU × account × event",
        "Group events into discount depth buckets (0–10%, 10–20%, 20–30%, 30%+)",
        "Separate response curves by channel — grocery and c-store show distinct lift profiles",
        "Separate response curves by pack size — 12-packs show a 30%+ discount cliff",
        "Identify the depth where marginal lift no longer justifies additional margin give-up",
    ],
    ml_outputs=[
        "Expected lift % at each discount depth for this SKU × channel × pack size",
        "Optimal discount depth recommendation — where lift peaks before the cliff",
        "Incremental vs. cannibalized demand breakdown at each discount level",
    ],
    live_outputs=[
        "Historical promo events — discount depth, timing, and observed lift per event",
        "Lift trend by event type — TPR vs. display vs. feature support",
        "% weeks on promo trend — is this SKU becoming promo-dependent?",
        "Channel comparison — lift curve across grocery / mass / c-store for this SKU",
    ],
    phase=2,
    aevah_raci="Lift curve modeling · Depth optimization scoring · Channel × pack separation · Publish to Mo UI",
    built_raci="SPINS export deposit (including promo flags) · Review lift projections · Confirm discount depth guidance with trade marketing team",
)

add_sim_slide(
    num_label="MODULE 04",
    name="Demand Velocity & Forecast",
    answer="What will this SKU sell in the next 13 weeks, by retailer account?",
    inputs=[
        ("52-week weekly units sold", "SPINS", "Demand history"),
        ("Distribution (stores selling)", "SPINS", "Distribution control"),
        ("Average retail price", "SPINS", "Price signal"),
        ("% weeks on promotion", "SPINS", "Promo flag"),
        ("Retail account", "SPINS", "Per-account model"),
        ("Week end date", "SPINS", "Seasonality features"),
    ],
    operations=[
        "Normalize to non-zero TDP weeks — eliminate false troughs from shelf-presence gaps",
        "Compute rolling averages (4-week, 13-week) to smooth noise while preserving trend",
        "Build 12-month seasonality index from category raw monthly averages",
        "Flag promo weeks so model learns baseline vs. promo-lifted demand separately",
        "Train separate forecast per SKU × retailer — Kroger velocity ≠ Target velocity",
        "Fall back to exponential smoothing (ETS) for SKUs with <8 weeks of history",
    ],
    ml_outputs=[
        "13-week demand forecast per SKU × retailer account",
        "Confidence bands (low / base / high scenario)",
        "Promo-uplift scenario — demand if a promotion is planned in the forecast window",
        "Forecast accuracy benchmark: 4% wMAPE vs. 7–10% industry baseline",
    ],
    live_outputs=[
        "52-week actual demand trend by SKU and retailer account",
        "Distribution (stores stocking) trend over the same period",
        "Category average velocity — how this SKU compares to 13-week category norm",
        "Seasonal index overlay — expected lift or drag on the forecast period",
    ],
    phase=2,
    aevah_raci="Seasonality index · Forecast model training · ETS fallback · Confidence bands · Publish 13-week projections to Mo UI",
    built_raci="SPINS export deposit · Connor validates 13-week forecasts against FP&A actuals · Share actuals data for backtesting and accuracy measurement",
)

add_sim_slide(
    num_label="MODULE 05",
    name="Distribution & Launch Monitoring",
    answer="Is this new product ramping as expected — in enough stores, with the right velocity?",
    inputs=[
        ("% of stores selling", "SPINS", "Distribution growth"),
        ("Total distribution points (TDP)", "SPINS", "Primary ramp metric"),
        ("Avg weekly units per store selling", "SPINS", "Velocity per store"),
        ("First week selling", "SPINS", "Launch anchor date"),
        ("Number of weeks selling", "SPINS", "Lifecycle position"),
    ],
    operations=[
        "Track weeks 1–16 of distribution build and velocity ramp for each new BUILT UPC",
        "Compare each SKU's ramp against category norms for the same pack type",
        "Classify new UPCs: new pack size, new flavor, or duplicate/relaunch",
        "Flag underperformance at 4, 8, and 13 weeks post-launch vs. expected ramp curve",
        "Suppress cannibalization scoring during ramp — avoids penalizing distribution-driven transfer",
    ],
    ml_outputs=[
        "Launch status at current week: On Track · Watch · Underperforming",
        "Expected ramp curve — where this SKU should be at weeks 8, 13, 16",
        "Predicted full-distribution velocity — expected units/store at maturity",
        "Cannibalization scoring readiness — flag for when SKU graduates to active scoring",
    ],
    live_outputs=[
        "Actual TDP ramp week by week since first week selling",
        "Actual velocity per store — is the product pulling through where stocked?",
        "Comparable launch benchmarks — similar BUILT or category launches in same retailer",
        "Data-maturity status — weeks until data gate clears for active scoring",
    ],
    phase=2,
    aevah_raci="Ramp monitoring model · Category benchmark comparisons · UPC classification gate · Launch status scoring · Publish to Mo UI",
    built_raci="SPINS export deposit · Confirm new UPC classifications (new SKU vs. relaunch) · Review launch status assessments with brand team",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MO TRENDS & EXTERNAL SIGNAL INTEGRATION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.333, 7.5, fill=OFF_WHITE)
rect(slide, 0, 0, 13.333, 1.1, fill=NAVY)
rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)

tb(slide, 0.5, 0.18, 10, 0.6,
   "Mo Trends & External Signal Integration", size=22, bold=True, color=WHITE)
tb(slide, 0.5, 0.72, 12, 0.32,
   "Competitive dashboard tiles showing velocity, price, distribution, and macro signals — all queried live, available from day one.",
   size=10.5, color=BLUE_LIGHT)

# ── Left: Mo Trends tiles ──
left_x = 0.4
tile_w = 5.9

rect(slide, left_x, 1.2, tile_w, 0.32, fill=NAVY)
tb(slide, left_x + 0.12, 1.24, tile_w - 0.24, 0.24,
   "06  ·  MO TRENDS — all data queried live, no ML scoring table",
   size=8, bold=True, color=BLUE_LIGHT)

tiles = [
    ("1 — Brand Velocity",   "Weekly units trend — 52-week history + 13-week rolling avg by product/retailer"),
    ("2 — Demand Forecast",  "13-week forward projection with confidence bands  [ML scoring table]"),
    ("3 — Pack Crossover",   "Single vs. multipack demand split across the BUILT product line"),
    ("4 — Price & Promo",    "ARP trend with promo event overlays and discount depth markers"),
    ("5 — Distribution Arc", "TDP build trajectory and ramp vs. category launch norms"),
    ("6 — Flavor Demand",    "Flavor-level velocity comparison across full BUILT flavor range"),
    ("7 — Macro Context",    "FRED GASDESW + UMCSENT — gas prices + U of Michigan consumer confidence"),
    ("8 — Velocity + Macro", "Stacked chart: brand velocity trend alongside macro signal overlay"),
]

tile_y = 1.52
row_h = 0.41
for i, (name, desc) in enumerate(tiles):
    bg = RGBColor(0xF8, 0xFA, 0xFD) if i % 2 == 0 else WHITE
    rect(slide, left_x, tile_y, tile_w, row_h, fill=bg, line=0.5,
         line_color=RGBColor(0xCC, 0xD8, 0xEE))
    tb(slide, left_x + 0.1, tile_y + 0.07, 2.15, row_h - 0.1,
       name, size=8.5, bold=True, color=DARK_TEXT)
    tb(slide, left_x + 2.3, tile_y + 0.07, tile_w - 2.4, row_h - 0.1,
       desc, size=8, color=MID_GREY)
    tile_y += row_h

# Filter bar note
rect(slide, left_x, tile_y + 0.08, tile_w, 0.36, fill=LIGHT_BLUE)
tb(slide, left_x + 0.12, tile_y + 0.14, tile_w - 0.24, 0.26,
   "Filter bar:  Product · Brand · Channel (FOOD / MASS / C-STORE) · Account · Units / $ toggle",
   size=8.5, color=DARK_TEXT)
tb(slide, left_x + 0.12, tile_y + 0.31, tile_w - 0.24, 0.16,
   "Demo defaults: BUILT Puff  +  KROGER  +  WALMART    |    MULO excluded from all Trends views",
   size=7.5, color=MID_GREY, italic=True)

# ── Right: External Signal Integration ──
right_x = 6.7
right_w  = 6.25

rect(slide, right_x, 1.2, right_w, 0.32, fill=GREEN)
tb(slide, right_x + 0.12, 1.24, right_w - 0.24, 0.24,
   "07  ·  EXTERNAL SIGNALS — commercial redistribution required for all",
   size=8, bold=True, color=WHITE)

# Live section
rect(slide, right_x, 1.52, right_w, 0.26, fill=GREEN_LIGHT)
tb(slide, right_x + 0.12, 1.55, right_w - 0.24, 0.2,
   "LIVE IN PRODUCTION", size=7.5, bold=True, color=GREEN)

live_sigs = [
    ("FRED — GASDESW",  "Weekly U.S. gas price index (St. Louis Fed)"),
    ("FRED — UMCSENT",  "U. of Michigan consumer confidence index"),
]
sig_y = 1.78
for name, desc in live_sigs:
    rect(slide, right_x, sig_y, right_w, 0.36, fill=WHITE, line=0.5,
         line_color=RGBColor(0xCC, 0xD8, 0xEE))
    rect(slide, right_x, sig_y, 0.07, 0.36, fill=GREEN)
    tb(slide, right_x + 0.16, sig_y + 0.06, 2.1, 0.26,
       name, size=9, bold=True, color=DARK_TEXT)
    tb(slide, right_x + 2.35, sig_y + 0.06, right_w - 2.45, 0.26,
       desc, size=8, color=MID_GREY)
    sig_y += 0.36

# Planned section
sig_y += 0.1
rect(slide, right_x, sig_y, right_w, 0.26, fill=AMBER_LIGHT)
tb(slide, right_x + 0.12, sig_y + 0.04, right_w - 0.24, 0.2,
   "PLANNED — OPEN ITEMS", size=7.5, bold=True, color=AMBER)
sig_y += 0.26

planned_sigs = [
    ("Open-Meteo Weather",    "Regional precip + temperature for seasonal demand context"),
    ("BLS CPI — Food",        "CPI food sub-indices for consumer price-pressure signal"),
    ("Kalshi Markets",        "Prediction market odds for macro events (recession, rates)"),
    ("Amazon BSR / Helium 10","Digital shelf velocity proxy for e-commerce demand"),
    ("EIA / USDA Commodities","Cocoa, whey, oat prices — ingredient cost-pressure signal"),
]
for name, desc in planned_sigs:
    rect(slide, right_x, sig_y, right_w, 0.36, fill=WHITE, line=0.5,
         line_color=RGBColor(0xCC, 0xD8, 0xEE))
    rect(slide, right_x, sig_y, 0.07, 0.36, fill=AMBER)
    tb(slide, right_x + 0.16, sig_y + 0.06, 2.1, 0.26,
       name, size=9, bold=True, color=DARK_TEXT)
    tb(slide, right_x + 2.35, sig_y + 0.06, right_w - 2.45, 0.26,
       desc, size=8, color=MID_GREY)
    sig_y += 0.36

footnote(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PHASE 1 PROJECT TIMELINE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.333, 7.5, fill=OFF_WHITE)
rect(slide, 0, 0, 13.333, 1.1, fill=NAVY)
rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)

tb(slide, 0.5, 0.18, 10, 0.6,
   "Phase 1 Project Timeline", size=24, bold=True, color=WHITE)
tb(slide, 0.5, 0.72, 12.5, 0.32,
   "Duration estimates are a starting point — to be validated with Brian before committing to dates. "
   "Durations reflect data readiness, not Aevah engineering capacity.",
   size=10.5, color=BLUE_LIGHT)

# Horizontal timeline arrow
arrow_y = 2.05
rect(slide, 0.5, arrow_y, 12.35, 0.06, fill=BLUE)
tb(slide, 12.7, arrow_y - 0.12, 0.4, 0.3, "▶", size=14, color=BLUE)

# Four milestone nodes
milestones = [
    ("Weeks 1–2",  "Data Commitment\n& Handoff",
     "BUILT confirms SPINS export format and schedule\n"
     "Aevah confirms storage, file-watch trigger, and QA gate configuration\n"
     "Both teams agree on data scope and contact points",
     BLUE, LIGHT_BLUE),
    ("Weeks 3–6",  "Data Onboarding\n& Validation",
     "Aevah: ingest → QA → feature engineering → model training → scoring\n"
     "Aevah delivers QA report for BUILT review\n"
     "BUILT validates: does the data look right?",
     GREEN, GREEN_LIGHT),
    ("Weeks 7–10", "Phase 1 Preview\n& Feedback",
     "Cannibalization Risk + Price Elasticity live in Mo UI\n"
     "Live demo with Brian + Connor — FOOD-channel accounts\n"
     "Connor + team provide UI feedback and projection validation",
     AMBER, AMBER_LIGHT),
    ("Weeks 11+",  "Iterate &\nPhase 2 Scope",
     "Address feedback; tune models based on BUILT market knowledge\n"
     "Phase 1 milestone sign-off by both teams\n"
     "Begin scoping Modules 3–5 and automation roadmap for Phase 2",
     PURPLE, PURPLE_LIGHT),
]

node_x_positions = [0.5, 3.75, 7.0, 10.25]
node_w = 2.85
node_h = 4.5
node_top = 1.35

for (week_label, title, body, header_color, body_color), nx in zip(milestones, node_x_positions):
    # Circle on timeline
    cx = nx + node_w / 2
    circle = slide.shapes.add_shape(9, Inches(cx - 0.16), Inches(arrow_y - 0.14),
                                     Inches(0.32), Inches(0.32))
    circle.fill.solid()
    circle.fill.fore_color.rgb = header_color
    circle.line.fill.background()

    # Card
    rect(slide, nx, node_top, node_w, node_h, fill=WHITE,
         line=0.75, line_color=RGBColor(0xC0, 0xCC, 0xE4))
    rect(slide, nx, node_top, node_w, 0.5, fill=header_color)
    tb(slide, nx + 0.12, node_top + 0.06, node_w - 0.24, 0.18,
       week_label, size=7.5, bold=True, color=WHITE)
    tb(slide, nx + 0.12, node_top + 0.24, node_w - 0.24, 0.3,
       title, size=11, bold=True, color=WHITE)

    # Body
    rect(slide, nx, node_top + 0.5, node_w, node_h - 0.5, fill=body_color,
         line=0.5, line_color=RGBColor(0xC0, 0xCC, 0xE4))
    tb(slide, nx + 0.14, node_top + 0.6, node_w - 0.28, node_h - 0.7,
       body, size=9, color=DARK_TEXT)

# Bottom note
rect(slide, 0.5, 6.18, 12.35, 0.52, fill=LIGHT_BLUE)
tb(slide, 0.65, 6.24, 12.1, 0.42,
   "Phase 2 (Modules 3–5: Promotional Response, Demand Forecast, Launch Monitoring) begins after Phase 1 sign-off milestone. "
   "Automation roadmap (incremental ingest, orchestration shell, file-watch trigger) runs in parallel as an Aevah-internal engineering sprint — no BUILT action required.",
   size=9, color=DARK_TEXT)

footnote(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — KICKOFF ALIGNMENT & NEXT STEPS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.333, 7.5, fill=NAVY_DEEP)
rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)
rect(slide, 0, 0, 13.333, 1.4, fill=RGBColor(0x14, 0x28, 0x50))

tb(slide, 0.5, 0.28, 10, 0.7,
   "Kickoff Alignment & Next Steps", size=26, bold=True, color=WHITE)
tb(slide, 0.5, 0.9, 12, 0.38,
   "The goal of this meeting: leave with a shared understanding of who owns what, "
   "what each team is committing to, and a validated path to Phase 1 delivery.",
   size=11, color=BLUE_LIGHT, italic=True)

# Two columns: What we need to align on (left) + What happens after (right)
left_x, left_w = 0.5, 5.8
right_x, right_w = 6.7, 6.2
col_top = 1.55

rect(slide, left_x, col_top, left_w, 0.32, fill=BLUE)
tb(slide, left_x + 0.14, col_top + 0.06, left_w - 0.28, 0.22,
   "WHAT WE NEED TO AGREE IN THIS MEETING", size=9, bold=True, color=WHITE)

align_items = [
    ("Data scope",
     "Which SPINS fields are we using? What supplemental data (ERP, DW) is in scope for Phase 1?"),
    ("Data handoff mechanics",
     "Who deposits the SPINS export? Where? On what cadence? Who confirms receipt?"),
    ("RACI sign-off",
     "Both teams confirm the responsibility matrix — who does what at each stage."),
    ("Timeline validation",
     "Are the week estimates realistic? Brian to flag any constraints (data availability, staff)."),
    ("Stakeholder map",
     "Who else needs to be in the loop? Are there voices we haven't heard from yet?"),
    ("Phase 1 success criteria",
     "What does a successful Phase 1 look like? How will we know we're ready for Phase 2?"),
]

item_y = col_top + 0.38
for label, desc in align_items:
    rect(slide, left_x + 0.12, item_y + 0.1, 0.08, 0.08, fill=BLUE)
    tb(slide, left_x + 0.28, item_y + 0.02, left_w - 0.38, 0.22,
       label, size=10, bold=True, color=WHITE)
    tb(slide, left_x + 0.28, item_y + 0.24, left_w - 0.38, 0.26,
       desc, size=9, color=BLUE_LIGHT)
    item_y += 0.6

rect(slide, right_x, col_top, right_w, 0.32, fill=GREEN)
tb(slide, right_x + 0.14, col_top + 0.06, right_w - 0.28, 0.22,
   "WHAT HAPPENS AFTER THIS MEETING", size=9, bold=True, color=WHITE)

next_items = [
    ("1",  "Data Commitment confirmed",
     "BUILT provides written confirmation of SPINS export format and deposit timeline"),
    ("2",  "Aevah configures ingest pipeline",
     "File-watch trigger set up; QA gates defined against agreed field list"),
    ("3",  "First data deposit",
     "BUILT deposits initial SPINS extract; Aevah runs ingest + QA; report shared with BUILT"),
    ("4",  "Feature engineering + model training",
     "Aevah runs Phase 1 models (cannibalization risk + price elasticity) — no BUILT action needed"),
    ("5",  "Phase 1 live demo",
     "Aevah hosts Mo walkthrough with Brian + Connor; projection validation session"),
    ("6",  "Feedback + iteration",
     "Connor's team reviews findings; UI adjustments made; Phase 2 scoping begins"),
]

item_y = col_top + 0.38
for num, title, desc in next_items:
    rect(slide, right_x + 0.12, item_y + 0.06, 0.26, 0.26,
         fill=RGBColor(0x20, 0x40, 0x70))
    tb(slide, right_x + 0.12, item_y + 0.07, 0.26, 0.22,
       num, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, right_x + 0.46, item_y + 0.02, right_w - 0.56, 0.22,
       title, size=10, bold=True, color=WHITE)
    tb(slide, right_x + 0.46, item_y + 0.24, right_w - 0.56, 0.22,
       desc, size=8.5, color=BLUE_LIGHT)
    item_y += 0.56

tb(slide, 0.5, 7.1, 12.5, 0.3,
   "Mo by Aevah  ·  BUILT Kickoff  ·  August 2026  ·  Confidential",
   size=8, color=MID_GREY, align=PP_ALIGN.RIGHT)


# ─── Save ────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "..", "mockups", "mo_data_ops_kickoff.pptx")
out_path = os.path.normpath(out_path)
prs.save(out_path)
print(f"Saved: {out_path}")
